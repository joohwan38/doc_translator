# pptx_handler.py
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN # PP_ALIGN may not be used directly, but good to have
from pptx.enum.lang import MSO_LANGUAGE_ID
import traceback

import os
import io
import logging
import re
# from datetime import datetime # Not used directly in this file in the latest version
from PIL import Image
# import tempfile # Not used directly in this file in the latest version
# import shutil # Not used directly in this file in the latest version

import config
from interfaces import AbsPptxProcessor, AbsTranslator, AbsOcrHandler, AbsOllamaService
from utils import setup_task_logging

from typing import Optional, Dict, Any, List, Tuple, Callable, TypedDict, IO


logger = logging.getLogger(__name__)

# Regular expression to find characters that are generally considered meaningful for translation.
# This helps in skipping purely symbolic or numeric text.
MEANINGFUL_CHAR_PATTERN = re.compile(
    r'[a-zA-Z'          # Basic Latin
    r'\u00C0-\u024F'    # Latin Extended-A and B (covers many European languages)
    r'\u1E00-\u1EFF'    # Latin Extended Additional
    # Add other scripts as needed by your target languages
    r'\u0600-\u06FF'    # Arabic
    r'\u0750-\u077F'    # Arabic Supplement
    r'\u08A0-\u08FF'    # Arabic Extended-A
    r'\u3040-\u30ff'    # Hiragana, Katakana (Japanese)
    r'\u3131-\uD79D'    # Hangul Compatibility Jamo, Hangul Syllables (Korean)
    r'\u4e00-\u9fff'    # CJK Unified Ideographs (Chinese, Japanese, Korean)
    r'\u0E00-\u0E7F'    # Thai
    r']'
)

def should_skip_translation(text: str) -> bool:
    """
    Determines if a given text string should be skipped for translation based on its content.
    Skips empty strings, strings with no meaningful characters, or strings with a low ratio of meaningful characters.
    """
    if not text:
        return True
    stripped_text = text.strip()
    if not stripped_text:
        logger.debug(f"Skipping translation (empty or whitespace only): '{text[:50]}...'")
        return True

    # Check if any meaningful character exists
    if not MEANINGFUL_CHAR_PATTERN.search(stripped_text):
        logger.debug(f"Skipping translation (no meaningful characters): '{stripped_text[:50]}...'")
        return True

    text_len = len(stripped_text)
    # For very short texts (e.g., <= 3 chars), translate if any meaningful char is present.
    if text_len <= 3:
        if MEANINGFUL_CHAR_PATTERN.search(stripped_text):
            logger.debug(f"Attempting translation (short string with meaningful char): '{stripped_text}'")
            return False
        else:
            logger.debug(f"Skipping translation (short string, no meaningful char): '{stripped_text}'")
            return True

    # For longer texts, check the ratio of meaningful characters.
    meaningful_chars_count = len(MEANINGFUL_CHAR_PATTERN.findall(stripped_text))
    ratio = meaningful_chars_count / text_len
    if ratio < config.MIN_MEANINGFUL_CHAR_RATIO_SKIP:
        logger.debug(f"Skipping translation (low meaningful char ratio {ratio:.2f}, threshold: {config.MIN_MEANINGFUL_CHAR_RATIO_SKIP}): '{stripped_text[:50]}...'")
        return True

    logger.debug(f"Attempting translation (conditions met): '{stripped_text[:50]}...'")
    return False


def is_ocr_text_valid(text: str) -> bool:
    """
    Determines if OCR-extracted text is valid enough for translation.
    Similar to should_skip_translation but might use different thresholds or criteria.
    """
    if not text:
        return False
    stripped_text = text.strip()
    if not stripped_text:
        return False

    if not MEANINGFUL_CHAR_PATTERN.search(stripped_text):
        logger.debug(f"OCR text validation skip (no meaningful characters): '{stripped_text[:50]}...'")
        return False

    text_len = len(stripped_text)
    if text_len <= 2: # Stricter for very short OCR text
        if MEANINGFUL_CHAR_PATTERN.search(stripped_text):
             logger.debug(f"OCR text valid (very short string with meaningful char): '{stripped_text}'")
             return True
        else:
            logger.debug(f"OCR text validation skip (very short, no meaningful char): '{stripped_text}'")
            return False

    meaningful_chars_count = len(MEANINGFUL_CHAR_PATTERN.findall(stripped_text))
    ratio = meaningful_chars_count / text_len
    if ratio < config.MIN_MEANINGFUL_CHAR_RATIO_OCR: # Potentially different threshold for OCR
        logger.debug(f"OCR text validation skip (low meaningful char ratio {ratio:.2f}, threshold: {config.MIN_MEANINGFUL_CHAR_RATIO_OCR}): '{stripped_text[:50]}...'")
        return False

    logger.debug(f"OCR text valid (conditions met): '{stripped_text[:50]}...'")
    return True

class TranslationJob(TypedDict):
    original_text: str
    context: Dict[str, Any] # Contains info like slide_idx, shape_obj_ref, item_type_internal, etc.
    is_ocr: bool
    char_count: int

class PptxHandler(AbsPptxProcessor):
    def __init__(self):
        pass

    def get_file_info(self, file_path: str) -> Dict[str, Any]: # Return type changed for error field
        """
        Analyzes a PPTX file and returns information about its content.
        """
        logger.info(f"Starting file info analysis for: {file_path}")
        info: Dict[str, Any] = { # Changed to Any to accommodate 'error' field
            "slide_count": 0,
            "text_elements_count": 0,       # Number of distinct shapes/table cells containing translatable text
            "total_text_char_count": 0, # Total characters in translatable text
            "image_elements_count": 0,      # Number of picture shapes
            "chart_elements_count": 0       # Number of chart shapes
        }
        try:
            prs = Presentation(file_path)
            info["slide_count"] = len(prs.slides)
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    is_text_element_counted_for_this_shape = False # Ensure a shape is counted once as a text element even if it has multiple paragraphs/runs
                    if shape.has_text_frame and hasattr(shape.text_frame, 'text') and \
                       shape.text_frame.text and shape.text_frame.text.strip():
                        current_shape_text = shape.text_frame.text
                        if not should_skip_translation(current_shape_text):
                            if not is_text_element_counted_for_this_shape:
                                info["text_elements_count"] += 1
                                is_text_element_counted_for_this_shape = True
                            info["total_text_char_count"] += len(current_shape_text) # Sum of all characters to translate

                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        info["image_elements_count"] += 1
                    elif shape.has_table:
                        table_text_element_count_incremented = False
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if hasattr(cell.text_frame, 'text') and cell.text_frame.text and cell.text_frame.text.strip():
                                    cell_text = cell.text_frame.text
                                    if not should_skip_translation(cell_text):
                                        if not table_text_element_count_incremented: # Count table once if it has any translatable cell
                                            info["text_elements_count"] += 1
                                            table_text_element_count_incremented = True
                                        info["total_text_char_count"] += len(cell_text)
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        info["chart_elements_count"] +=1
            logger.info(
                f"File analysis complete: Slides: {info['slide_count']}, "
                f"TextElements: {info['text_elements_count']} (TotalChars: {info['total_text_char_count']}), "
                f"Images: {info['image_elements_count']}, Charts: {info['chart_elements_count']}"
            )
        except Exception as e:
            logger.error(f"Error analyzing file info for '{os.path.basename(file_path)}': {e}", exc_info=True)
            info = { "slide_count": -1, "text_elements_count": -1, "total_text_char_count": -1,
                     "image_elements_count": -1, "chart_elements_count": -1, "error": str(e)}
        return info

    def _get_style_properties(self, font_obj: Any) -> Dict[str, Any]:
        """Extracts common style properties from a font object."""
        if font_obj is None:
            return {}
        style_props: Dict[str, Any] = {
            'name': None, 'size': None, 'bold': None, 'italic': None,
            'underline': None, 'color_rgb': None, 'color_theme_index': None,
            'color_brightness': 0.0, 'language_id': None # MSO_LANGUAGE_ID.NONE (0) or other
        }
        # Safe attribute access
        for prop, attr_name in [('name', 'name'), ('size', 'size'), ('bold', 'bold'),
                                ('italic', 'italic'), ('underline', 'underline')]:
            try: style_props[prop] = getattr(font_obj, attr_name, None)
            except AttributeError: pass # Some attributes might not exist on all font-like objects

        try: style_props['language_id'] = font_obj.language_id
        except (ValueError, AttributeError): pass # Handles invalid language IDs or missing attribute

        # Color extraction
        if hasattr(font_obj, 'color') and hasattr(font_obj.color, 'type'):
            color_format = font_obj.color
            if color_format.type == MSO_COLOR_TYPE.RGB:
                try: style_props['color_rgb'] = RGBColor(*color_format.rgb) # Store as RGBColor object
                except AttributeError: pass
            elif color_format.type == MSO_COLOR_TYPE.SCHEME:
                try: style_props['color_theme_index'] = color_format.theme_color
                except AttributeError: pass
                try: style_props['color_brightness'] = color_format.brightness
                except AttributeError: style_props['color_brightness'] = 0.0 # Default brightness
        return style_props

    def _get_text_style(self, run_obj: Any) -> Dict[str, Any]:
        """Extracts style from a run, including hyperlink."""
        style = self._get_style_properties(run_obj.font)
        try:
            style['hyperlink_address'] = run_obj.hyperlink.address if run_obj.hyperlink and run_obj.hyperlink.address else None
        except AttributeError:
            style['hyperlink_address'] = None # No hyperlink or not accessible
        return style

    def _apply_style_properties(self, target_font_obj: Any, style_to_apply: Dict[str, Any]):
        """Applies extracted style properties to a target font object."""
        if not style_to_apply or target_font_obj is None:
            return

        font = target_font_obj
        # Apply basic font properties
        if style_to_apply.get('name') is not None: font.name = style_to_apply['name']
        if style_to_apply.get('size') is not None: font.size = style_to_apply['size']
        if style_to_apply.get('bold') is not None: font.bold = style_to_apply['bold']
        if style_to_apply.get('italic') is not None: font.italic = style_to_apply['italic']
        if style_to_apply.get('underline') is not None: font.underline = style_to_apply['underline']

        # Apply color
        applied_color = False
        if style_to_apply.get('color_rgb') is not None:
            try:
                # If stored as RGBColor object, use it directly. If tuple, convert.
                rgb_color_val = style_to_apply['color_rgb']
                if isinstance(rgb_color_val, RGBColor):
                    font.color.rgb = rgb_color_val
                    applied_color = True
                elif isinstance(rgb_color_val, tuple) and len(rgb_color_val) == 3:
                    font.color.rgb = RGBColor(*rgb_color_val)
                    applied_color = True
            except Exception as e:
                logger.warning(f"Failed to apply RGB color {style_to_apply['color_rgb']}: {e}")

        if not applied_color and style_to_apply.get('color_theme_index') is not None:
            try:
                theme_color_val = style_to_apply['color_theme_index']
                # Ensure it's a valid MSO_THEME_COLOR_INDEX enum member
                if isinstance(theme_color_val, MSO_THEME_COLOR_INDEX):
                    font.color.theme_color = theme_color_val
                # If it's an int, try to convert to enum (python-pptx might handle this)
                elif isinstance(theme_color_val, int):
                     font.color.theme_color = MSO_THEME_COLOR_INDEX(theme_color_val) # May raise ValueError if invalid
                else:
                    logger.warning(f"Invalid theme color index type/value: {theme_color_val}")

                brightness_val = float(style_to_apply.get('color_brightness', 0.0))
                font.color.brightness = max(-1.0, min(1.0, brightness_val)) # Clamp between -1 and 1
            except Exception as e:
                logger.warning(f"Failed to apply theme color {style_to_apply.get('color_theme_index')}: {e}")
        
        # Apply language ID
        if style_to_apply.get('language_id') is not None:
            try:
                lang_id_val = style_to_apply['language_id']
                if isinstance(lang_id_val, MSO_LANGUAGE_ID):
                    font.language_id = lang_id_val
                elif isinstance(lang_id_val, int) : # check if it's a valid MSO_LANGUAGE_ID value
                    try:
                        font.language_id = MSO_LANGUAGE_ID(lang_id_val)
                    except ValueError: # python-pptx raises ValueError for invalid enum values
                        logger.warning(f"Invalid MSO_LANGUAGE_ID value: {lang_id_val}. Using default.")
                # else: logger.debug(f"Language ID '{lang_id_val}' type ({type(lang_id_val)}) not directly applicable. Ignored.")
            except Exception as e_lang:
                logger.warning(f"Failed to apply language_id '{style_to_apply['language_id']}': {e_lang}")

    def _apply_text_style(self, run_obj: Any, style_to_apply: Dict[str, Any]):
        """Applies style to a run object, including font and hyperlink."""
        if not style_to_apply or run_obj is None:
            return
        self._apply_style_properties(run_obj.font, style_to_apply)
        
        # Apply hyperlink (if any)
        # Note: Clearing existing hyperlink before applying a new one might be needed if overwriting.
        # For simplicity, this example applies if address is present.
        # If 'hyperlink_address' is None in style_to_apply, it implies no hyperlink or removal.
        hlink_address = style_to_apply.get('hyperlink_address')
        if hlink_address:
            try:
                run_obj.hyperlink.address = hlink_address
            except Exception as e_hlink:
                logger.warning(f"Failed to apply hyperlink '{hlink_address}': {e_hlink}")
        # To remove a hyperlink, you might need to do run.hyperlink.address = None, or more involved XML manipulation
        # For now, we are only restoring, so if it was None, it remains None unless explicitly set.

    def _apply_paragraph_style(self, paragraph_obj: Any, para_style_info: Dict[str, Any]):
        """Applies paragraph-level styles (alignment, level, spacing)."""
        if not para_style_info or paragraph_obj is None:
            return
        try:
            if para_style_info.get('alignment') is not None:
                paragraph_obj.alignment = para_style_info['alignment']
            paragraph_obj.level = para_style_info.get('level', 0) # Default to level 0

            # Apply spacing properties if they exist and are valid (e.g., Pt objects or None)
            for space_attr in ['space_before', 'space_after', 'line_spacing']:
                if para_style_info.get(space_attr) is not None:
                    try:
                        setattr(paragraph_obj, space_attr, para_style_info[space_attr])
                    except Exception as e_space:
                         logger.debug(f"Could not set paragraph spacing '{space_attr}' to '{para_style_info[space_attr]}': {e_space}")
        except Exception as e:
            logger.warning(f"Error applying paragraph style: {e}")

    def _safe_clear_text_frame(self, text_frame_obj: Any, log_func: Optional[Callable[[str], None]] = None):
        """Safely clears all paragraphs and runs from a text frame."""
        try:
            # Standard way to clear: remove all paragraphs
            # Iterating and removing paragraphs one by one can be problematic if the collection changes.
            # A common approach is to remove them in reverse or clear text from each paragraph.
            # text_frame.clear() is the most direct if available and works as expected.
            
            # Check if text_frame.clear() method exists and use it
            if hasattr(text_frame_obj, 'clear') and callable(text_frame_obj.clear):
                 text_frame_obj.clear()
                 if log_func: log_func("      Text frame cleared using .clear() method.")
                 return

            # Fallback: If .clear() is not available or to be absolutely sure, manipulate paragraphs directly
            # This part might be redundant if .clear() works well.
            # Removing paragraphs from _txBody (XML element) is more low-level.
            if hasattr(text_frame_obj, '_element') and text_frame_obj._element is not None:
                txBody = text_frame_obj._element # This is txBody
                # Get all paragraph child elements ('a:p')
                paragraphs_xml_to_remove = [child for child in list(txBody) if child.tag.endswith('}p')]
                if paragraphs_xml_to_remove:
                    if log_func: log_func(f"      Attempting to remove {len(paragraphs_xml_to_remove)} existing paragraphs at XML level...")
                    for p_xml in paragraphs_xml_to_remove:
                        try:
                            txBody.remove(p_xml)
                        except Exception as e_xml_remove:
                            logger.debug(f"Ignored error removing XML paragraph: {e_xml_remove}")
                    if log_func: log_func(f"      XML level paragraph removal complete.")
                else:
                    if log_func: log_func(f"      No existing XML paragraphs found to remove.")
            else: # If no _element, try to clear paragraphs individually
                for p in reversed(text_frame_obj.paragraphs): # Iterate in reverse for safe removal
                    p.clear() # Clear content of paragraph
                    # To remove the paragraph element itself (p_elem = p._p; p_elem.getparent().remove(p_elem)) - more complex
                if log_func: log_func("      Text frame paragraphs cleared individually (fallback).")

        except Exception as e:
            logger.warning(f"Error clearing text frame: {e}", exc_info=True)
            if log_func: log_func(f"      Error clearing text frame: {e}")


    def _apply_translated_text_to_frame(self, text_frame_obj: Any, translated_text: str,
                                      original_para_styles_list: List[Dict[str, Any]],
                                      item_name_for_log: str, log_func: Optional[Callable[[str], None]]):
        # 1. Backup original text frame properties (기존과 동일)
        original_tf_properties = {}
        tf_props_to_backup = ['auto_size', 'vertical_anchor', 'word_wrap',
                              'margin_left', 'margin_right', 'margin_top', 'margin_bottom']
        for prop_name in tf_props_to_backup:
            try:
                original_tf_properties[prop_name] = getattr(text_frame_obj, prop_name, None)
            except AttributeError:
                original_tf_properties[prop_name] = None
                logger.debug(f"Attribute '{prop_name}' not found on text_frame for '{item_name_for_log}'.")

        # 2. Prepare text frame for new content (기존과 동일)
        try:
            if original_tf_properties.get('auto_size') is not None and \
               original_tf_properties['auto_size'] != MSO_AUTO_SIZE.NONE:
                text_frame_obj.auto_size = MSO_AUTO_SIZE.NONE
            if original_tf_properties.get('word_wrap') is False :
                 text_frame_obj.word_wrap = True
        except Exception as e_prop_set:
            logger.debug(f"Error setting temporary text_frame properties for '{item_name_for_log}': {e_prop_set}")

        # 3. Clear existing content (기존과 동일)
        self._safe_clear_text_frame(text_frame_obj, log_func)
        # text_frame_obj.clear() 호출 후, text_frame_obj.paragraphs 에는 비어있는 단락이 하나 있을 가능성이 높음

        # 4. Add translated text with styles
        try:
            # 번역된 텍스트가 비어있거나 공백만 있는 경우 처리
            if not translated_text.strip():
                p_target = None
                if text_frame_obj.paragraphs: # clear() 후 남아있는 단락 사용
                    p_target = text_frame_obj.paragraphs[0]
                    p_target.clear() # 기존 단락의 내용(런)을 모두 지움
                else: # 만약 단락이 없다면 (이론상 발생 안 함)
                    p_target = text_frame_obj.add_paragraph()

                # 원본 첫 번째 단락 스타일 적용 시도
                if original_para_styles_list:
                    first_para_style_info = original_para_styles_list[0]
                    self._apply_paragraph_style(p_target, first_para_style_info)
                    run_style_to_apply = first_para_style_info.get('runs')[0].get('style', {}) if first_para_style_info.get('runs') else \
                                         first_para_style_info.get('paragraph_default_run_style', {})
                    run_new = p_target.add_run()
                    run_new.text = " " # 공백 한 칸 삽입
                    self._apply_text_style(run_new, run_style_to_apply)
                else: # 스타일 정보 없을 경우
                    run_new = p_target.add_run()
                    run_new.text = " "
                if log_func: log_func(f"      Applied empty (space) translated text to '{item_name_for_log}'.")
            else:
                # 번역된 텍스트 내용이 있는 경우
                lines = translated_text.splitlines()
                if not lines and translated_text: # splitlines가 비었지만 원본 텍스트는 있는 경우 (개행 없는 한 줄)
                    lines = [translated_text]

                # --- 선행 줄바꿈 문자 유물 감지 로직 (기존과 동일) ---
                is_leading_newline_artifact = False
                if len(lines) > 0 and not lines[0].strip() and \
                   any(line.strip() for line in lines[1:]):
                    is_leading_newline_artifact = True
                    if log_func: log_func(f"      Detected and will skip artifact leading empty line for '{item_name_for_log}'.")
                # --- 로직 끝 ---

                first_line_written = False
                paragraphs_added_count = 0 # 실제 내용이 있는 단락 수
                
                for i, line_text in enumerate(lines):
                    if is_leading_newline_artifact and i == 0: # 유물 줄바꿈 건너뛰기
                        continue

                    p_current = None
                    # 첫 번째 줄을 쓸 차례이고, text_frame에 이미 단락이 있다면 (clear()로 인해 남은 단락) 그것을 사용
                    if not first_line_written and text_frame_obj.paragraphs:
                        p_current = text_frame_obj.paragraphs[0]
                        p_current.clear() # 기존 단락의 내용(런)을 모두 지움
                    else:
                        # 이후 줄이거나, 처음부터 단락이 없었다면 새로 추가
                        p_current = text_frame_obj.add_paragraph()
                    
                    paragraphs_added_count +=1

                    current_line_text_to_write = line_text if line_text.strip() else " "

                    # 스타일 적용 로직 (기존과 유사하게, 현재 라인 인덱스와 유물 줄바꿈 고려)
                    style_application_line_index = i - (1 if is_leading_newline_artifact else 0)
                    para_style_to_apply = {}
                    run_style_to_apply = {}

                    if original_para_styles_list and style_application_line_index >= 0:
                        style_template_index = min(style_application_line_index, len(original_para_styles_list) - 1)
                        if style_template_index < len(original_para_styles_list): # 유효 인덱스 확인
                            para_style_info = original_para_styles_list[style_template_index]
                            para_style_to_apply = para_style_info
                            if para_style_info.get('runs') and para_style_info['runs']:
                                run_style_to_apply = para_style_info['runs'][0].get('style', {})
                            elif para_style_info.get('paragraph_default_run_style'):
                                run_style_to_apply = para_style_info['paragraph_default_run_style']
                    
                    self._apply_paragraph_style(p_current, para_style_to_apply)
                    run_new = p_current.add_run()
                    run_new.text = current_line_text_to_write
                    self._apply_text_style(run_new, run_style_to_apply)
                    
                    first_line_written = True

                # 모든 줄이 비어있거나 유물 처리되어 건너뛰어졌지만, 원본 번역에 내용이 있었던 경우
                if paragraphs_added_count == 0 and translated_text.strip():
                    if log_func: log_func(f"      All lines were skipped or empty for '{item_name_for_log}' but original translation had content. Adding default paragraph.")
                    p_fallback = None
                    if text_frame_obj.paragraphs: # 첫번째 단락을 재활용 시도
                        p_fallback = text_frame_obj.paragraphs[0]
                        p_fallback.clear()
                    else: # 없으면 새로 추가
                        p_fallback = text_frame_obj.add_paragraph()
                    
                    run_fallback = p_fallback.add_run()
                    run_fallback.text = translated_text.strip() # 최소한의 내용이라도 표시

                    if original_para_styles_list: # 첫번째 스타일이라도 적용 시도
                        first_para_style_info = original_para_styles_list[0]
                        self._apply_paragraph_style(p_fallback, first_para_style_info)
                        run_style_to_apply_fallback = first_para_style_info.get('runs')[0].get('style', {}) if first_para_style_info.get('runs') else \
                                                      first_para_style_info.get('paragraph_default_run_style', {})
                        self._apply_text_style(run_fallback, run_style_to_apply_fallback)

                if log_func: log_func(f"      Applied translated text to '{item_name_for_log}'.")
        except Exception as e_apply:
            logger.error(f"Error applying translated text to '{item_name_for_log}': {e_apply}", exc_info=True)
            if log_func: log_func(f"      ERROR: Failed to apply translated text to '{item_name_for_log}': {e_apply}")

        # 5. Restore original text frame properties
        try:
            for prop_name, original_value in original_tf_properties.items():
                if original_value is not None :
                    if prop_name == 'word_wrap' and text_frame_obj.word_wrap is True and original_value is False:
                        if log_func: log_func(f"      Keeping word_wrap=True for '{item_name_for_log}' for better i18n display (original was False).")
                        continue
                    setattr(text_frame_obj, prop_name, original_value)
        except Exception as e_prop_restore:
            logger.debug(f"Error restoring text_frame properties for '{item_name_for_log}': {e_prop_restore}")


    def translate_presentation_stage1(self, prs: Presentation, src_lang_ui_name: str, tgt_lang_ui_name: str,
                                        translator: AbsTranslator, ocr_handler: Optional[AbsOcrHandler],
                                        model_name: str, ollama_service: AbsOllamaService,
                                        font_code_for_render: str, task_log_filepath: str,
                                        progress_callback_item_completed: Optional[Callable[[Any, str, float, str], None]] = None,
                                        stop_event: Optional[Any] = None,
                                        image_translation_enabled: bool = True,
                                        ocr_temperature: Optional[float] = None
                                        ) -> bool:
        initial_log_lines_s1 = ["--- 1단계: 차트 외 요소 번역 시작 (PptxHandler - Stage 1) ---"]
        f_task_log_s1: Optional[IO[str]] = None
        log_func_s1: Optional[Callable[[str], None]] = None

        if task_log_filepath:
            f_task_log_s1, log_func_s1_temp = setup_task_logging(task_log_filepath, initial_log_lines_s1)
            if log_func_s1_temp:
                log_func_s1 = log_func_s1_temp
        
        main_log_prefix = "PPTX Stage 1:"
        logger.info(f"{main_log_prefix} Starting text and image (OCR) content collection.")
        if log_func_s1: log_func_s1("1단계: 텍스트 및 이미지(OCR) 내용 수집 중...")

        translation_jobs: List[TranslationJob] = []
        original_paragraph_styles_map: Dict[Tuple[int, Any, Any], List[Dict[str, Any]]] = {}

        try: 
            for slide_idx, slide in enumerate(prs.slides):
                if stop_event and stop_event.is_set():
                    logger.info(f"{main_log_prefix} Stop event detected during slide iteration for job collection.")
                    if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                    return False
                for shape_idx, shape in enumerate(slide.shapes): 
                    if stop_event and stop_event.is_set(): 
                        logger.info(f"{main_log_prefix} Stop event detected during shape iteration for job collection.")
                        if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                        return False

                    shape_id_for_key = getattr(shape, 'shape_id', f"s{slide_idx}_auto_idx{shape_idx}") 
                    element_name_for_log = shape.name or f"Slide{slide_idx+1}_Shape{shape_idx}(ID:{shape_id_for_key})"
                    item_base_context = {'slide_idx': slide_idx, 'shape_obj_ref': shape, 'name': element_name_for_log, 'shape_id_log': shape_id_for_key}
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        continue

                    if shape.has_text_frame and getattr(shape.text_frame, 'text', None) and shape.text_frame.text.strip():
                        original_text = shape.text_frame.text
                        if not should_skip_translation(original_text):
                            style_key = (slide_idx, shape_id_for_key, 'shape_text')
                            translation_jobs.append({
                                'original_text': original_text,
                                'context': {**item_base_context, 'item_type_internal': 'text_shape', 'style_unique_key': style_key},
                                'is_ocr': False, 'char_count': len(original_text)
                            })
                    elif shape.has_table:
                        for r_idx, row in enumerate(shape.table.rows):
                            for c_idx, cell in enumerate(row.cells):
                                if getattr(cell.text_frame, 'text', None) and cell.text_frame.text.strip():
                                    original_text = cell.text_frame.text
                                    if not should_skip_translation(original_text):
                                        style_key = (slide_idx, shape_id_for_key, ('table_cell', r_idx, c_idx))
                                        cell_log_name = f"{element_name_for_log}_R{r_idx}C{c_idx}"
                                        translation_jobs.append({
                                            'original_text': original_text,
                                            'context': {**item_base_context, 'name': cell_log_name, 
                                                        'item_type_internal': 'table_cell', 'row_idx': r_idx, 
                                                        'col_idx': c_idx, 'style_unique_key': style_key},
                                            'is_ocr': False, 'char_count': len(original_text)
                                        })
            
            if log_func_s1: log_func_s1(f"1단계 텍스트/표 번역 작업 수집 완료: {len(translation_jobs)}개 항목.")
            
            if not translation_jobs and not (image_translation_enabled and ocr_handler):
                is_any_chart_present = any(s.shape_type == MSO_SHAPE_TYPE.CHART for slide_obj_check in prs.slides for s in slide_obj_check.shapes)
                if not is_any_chart_present:
                    if log_func_s1: log_func_s1("1단계: 처리 대상(텍스트/표/OCR/차트) 없음. 건너뜀.")
                    if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                    return True

            texts_for_batch_translation = [job['original_text'] for job in translation_jobs if not job['is_ocr']]
            translated_texts_batch: List[str] = []

            if texts_for_batch_translation:
                if log_func_s1: log_func_s1(f"일반 텍스트 {len(texts_for_batch_translation)}개 일괄 번역 시작...")
                translated_texts_batch = translator.translate_texts_batch(
                    texts_for_batch_translation, src_lang_ui_name, tgt_lang_ui_name,
                    model_name, ollama_service, is_ocr_text=False, stop_event=stop_event
                )
                if stop_event and stop_event.is_set():
                    if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                    return False
                if len(texts_for_batch_translation) != len(translated_texts_batch):
                    logger.warning(f"{main_log_prefix} 텍스트 수와 번역 결과 수 불일치!")
                    if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                    return False
            
            current_batch_text_idx = 0
            for job_data in translation_jobs:
                if stop_event and stop_event.is_set(): break
                if job_data['is_ocr']: continue

                context = job_data['context']
                slide_idx_job = context['slide_idx']
                item_name_for_log_job = context['name']
                item_type_internal_job = context['item_type_internal']
                ui_feedback_task_type = "텍스트/표 적용"
                if item_type_internal_job == 'text_shape': ui_feedback_task_type = "텍스트 요소 적용"
                elif item_type_internal_job == 'table_cell': ui_feedback_task_type = "표 셀 내용 적용"
                
                text_frame_to_process: Optional[Any] = None
                if item_type_internal_job == 'text_shape':
                    shape_obj_ref = context.get('shape_obj_ref')
                    if shape_obj_ref and shape_obj_ref.has_text_frame:
                        text_frame_to_process = shape_obj_ref.text_frame
                elif item_type_internal_job == 'table_cell':
                    table_shape_obj_ref = context.get('shape_obj_ref') 
                    r_idx_job, c_idx_job = context['row_idx'], context['col_idx']
                    if table_shape_obj_ref and table_shape_obj_ref.has_table:
                        try:
                            text_frame_to_process = table_shape_obj_ref.table.cell(r_idx_job, c_idx_job).text_frame
                        except IndexError:
                            logger.error(f"{main_log_prefix} 테이블 셀 접근 오류: '{item_name_for_log_job}' at R{r_idx_job}C{c_idx_job}.")
                            if progress_callback_item_completed:
                                progress_callback_item_completed(f"슬라이드 {slide_idx_job + 1} 오류", "테이블 셀 접근 실패", float(job_data['char_count'] * config.WEIGHT_TEXT_CHAR), f"셀 ({r_idx_job},{c_idx_job})")
                            continue 

                if text_frame_to_process:
                    style_key_job = context['style_unique_key']
                    if style_key_job not in original_paragraph_styles_map:
                        collected_para_styles: List[Dict[str, Any]] = []
                        try:
                            for para in text_frame_to_process.paragraphs:
                                para_default_style = self._get_style_properties(para.font)
                                runs_info_list: List[Dict[str, Any]] = []
                                if para.runs:
                                    for run in para.runs:
                                        runs_info_list.append({'text': run.text, 'style': self._get_text_style(run)})
                                elif para.text and para.text.strip(): 
                                    runs_info_list.append({'text': para.text, 'style': self._get_text_style(para.runs[0] if para.runs else para)}) 
                                collected_para_styles.append({
                                    'runs': runs_info_list,
                                    'alignment': para.alignment, 'level': para.level,
                                    'space_before': para.space_before, 'space_after': para.space_after,
                                    'line_spacing': para.line_spacing,
                                    'paragraph_default_run_style': para_default_style
                                })
                            original_paragraph_styles_map[style_key_job] = collected_para_styles
                        except Exception as e_style_collect:
                            logger.error(f"{main_log_prefix} Error collecting styles for '{item_name_for_log_job}': {e_style_collect}", exc_info=True)
                            original_paragraph_styles_map[style_key_job] = [] 

                    translated_text_for_job = translated_texts_batch[current_batch_text_idx] if current_batch_text_idx < len(translated_texts_batch) else job_data['original_text']
                    current_batch_text_idx +=1
                    
                    if "오류:" not in translated_text_for_job: 
                        self._apply_translated_text_to_frame(
                            text_frame_to_process,
                            translated_text_for_job,
                            original_paragraph_styles_map.get(style_key_job, []), 
                            item_name_for_log_job,
                            log_func_s1
                        )
                
                if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                    progress_text_snippet = translated_text_for_job.strip().replace(chr(10), ' ')[:30] if "오류:" not in translated_text_for_job else job_data['original_text'].strip().replace(chr(10), ' ')[:30]
                    progress_callback_item_completed(
                        f"슬라이드 {slide_idx_job + 1}",
                        ui_feedback_task_type,
                        float(job_data['char_count'] * config.WEIGHT_TEXT_CHAR), 
                        progress_text_snippet
                    )
            
            if stop_event and stop_event.is_set(): # 일괄 번역 적용 후 중단 체크
                if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                return False

            # MODIFIED PART STARTS HERE
            if image_translation_enabled:
                if ocr_handler:
                    if log_func_s1: log_func_s1("\n--- 이미지 OCR 및 번역 적용 시작 (슬라이드별) ---")
                    logger.info(f"{main_log_prefix} Starting OCR and translation for images (per slide).")

                    for slide_idx_ocr, current_slide_for_ocr in enumerate(prs.slides):
                        if stop_event and stop_event.is_set():
                            if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                            return False
                        
                        if log_func_s1:
                            log_func_s1(f"  [S{slide_idx_ocr+1}] 이미지 OCR 분석 중. 총 도형 수: {len(current_slide_for_ocr.shapes)}")
                            for shape_idx_debug, shape_debug in enumerate(current_slide_for_ocr.shapes):
                                try:
                                    shape_type_str = shape_debug.shape_type.name if hasattr(shape_debug.shape_type, 'name') else str(shape_debug.shape_type)
                                    log_func_s1(f"    도형 {shape_idx_debug+1}: 타입='{shape_type_str}', 이름='{shape_debug.name if hasattr(shape_debug, 'name') else 'N/A'}'")
                                except Exception:
                                    log_func_s1(f"    도형 {shape_idx_debug+1}: 타입 정보 읽기 오류")
                        
                        picture_shapes_on_slide = [
                            s for s in current_slide_for_ocr.shapes
                            if hasattr(s, 'shape_type') and s.shape_type == MSO_SHAPE_TYPE.PICTURE
                        ]

                        if not picture_shapes_on_slide:
                            if log_func_s1: log_func_s1(f"  [S{slide_idx_ocr+1}] OCR 대상 이미지 없음.")
                        else:
                            if log_func_s1: log_func_s1(f"  [S{slide_idx_ocr+1}] {len(picture_shapes_on_slide)}개의 OCR 대상 이미지 발견.")

                            for shape_ocr_idx, shape_ocr in enumerate(picture_shapes_on_slide):
                                if stop_event and stop_event.is_set():
                                    if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                                    return False

                                item_name_ocr_log = shape_ocr.name or f"Slide{slide_idx_ocr+1}_Image{shape_ocr_idx+1}"
                                ocr_feedback_location = f"슬라이드 {slide_idx_ocr + 1}"
                                ocr_item_total_weight = float(config.WEIGHT_IMAGE)
                                current_ocr_item_processed_weight = 0.0
                                
                                if log_func_s1: log_func_s1(f"\n  [S{slide_idx_ocr+1}] OCR 처리 시도: '{item_name_ocr_log}'")
                                
                                try:
                                    img_bytes = shape_ocr.image.blob
                                    img_pil_original = Image.open(io.BytesIO(img_bytes))
                                    img_pil_rgb = img_pil_original.convert("RGB")
                                    original_img_format = img_pil_original.format
                                    
                                    ocr_detection_weight_portion = ocr_item_total_weight * 0.2
                                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                        progress_callback_item_completed(ocr_feedback_location, "ocr_status_detection_start", 0, f"'{item_name_ocr_log}' 감지 시작")

                                    ocr_results = ocr_handler.ocr_image(img_pil_rgb)
                                    if stop_event and stop_event.is_set():
                                        if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                                        return False
                                    
                                    ocr_result_count = len(ocr_results) if ocr_results else 0
                                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                        progress_callback_item_completed(
                                            ocr_feedback_location, 
                                            "ocr_status_detection_complete" if ocr_results else "ocr_status_detection_no_text",
                                            ocr_detection_weight_portion, 
                                            f"'{item_name_ocr_log}' ({ocr_result_count} 블록)"
                                        )
                                    current_ocr_item_processed_weight += ocr_detection_weight_portion
                                    if log_func_s1: log_func_s1(f"        '{item_name_ocr_log}' OCR 분석 완료. {ocr_result_count}개 블록 발견.")

                                    ocr_texts_for_translation: List[str] = []
                                    ocr_contexts_for_render: List[Dict[str, Any]] = []
                                    if ocr_results:
                                        for res_item in ocr_results:
                                            if not (isinstance(res_item, (list, tuple)) and len(res_item) >= 2): continue
                                            box_coords, text_conf_pair = res_item[0], res_item[1]
                                            text_angle = res_item[2] if len(res_item) > 2 else None
                                            if not (isinstance(text_conf_pair, (list, tuple)) and len(text_conf_pair) == 2): continue
                                            original_ocr_text, _ = text_conf_pair
                                            if is_ocr_text_valid(original_ocr_text) and not should_skip_translation(original_ocr_text):
                                                ocr_texts_for_translation.append(original_ocr_text)
                                                ocr_contexts_for_render.append({'box': box_coords, 'original_text': original_ocr_text, 'angle': text_angle})
                                    
                                    if ocr_texts_for_translation:
                                        ocr_translation_weight_portion = ocr_item_total_weight * 0.4
                                        if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                            progress_callback_item_completed(ocr_feedback_location, "ocr_status_translating_texts_start", 0, f"'{item_name_ocr_log}' ({len(ocr_texts_for_translation)}개) 번역 시작")
                                        
                                        translated_ocr_texts = translator.translate_texts_batch(
                                            ocr_texts_for_translation, src_lang_ui_name, tgt_lang_ui_name,
                                            model_name, ollama_service, is_ocr_text=True,
                                            ocr_temperature=ocr_temperature, stop_event=stop_event
                                        )
                                        if stop_event and stop_event.is_set():
                                            if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                                            return False
                                        
                                        if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                            progress_callback_item_completed(
                                                ocr_feedback_location, "ocr_status_translating_texts_complete",
                                                ocr_translation_weight_portion, f"'{item_name_ocr_log}' ({len(translated_ocr_texts)}개 결과)"
                                            )
                                        current_ocr_item_processed_weight += ocr_translation_weight_portion

                                        if len(ocr_texts_for_translation) == len(translated_ocr_texts):
                                            img_to_render_on = img_pil_original.copy()
                                            any_text_rendered_on_image = False
                                            ocr_rendering_total_alloc_weight = ocr_item_total_weight * 0.4
                                            
                                            if translated_ocr_texts:
                                                render_per_block_weight = ocr_rendering_total_alloc_weight / len(translated_ocr_texts) if translated_ocr_texts else 0
                                                for i, translated_text_render in enumerate(translated_ocr_texts):
                                                    if stop_event and stop_event.is_set(): break
                                                    render_ctx = ocr_contexts_for_render[i]
                                                    if "오류:" not in translated_text_render and translated_text_render.strip():
                                                        try:
                                                            img_to_render_on = ocr_handler.render_translated_text_on_image(
                                                                img_to_render_on, render_ctx['box'], translated_text_render,
                                                                font_code_for_render, render_ctx['original_text'], render_ctx['angle']
                                                            )
                                                            any_text_rendered_on_image = True
                                                        except Exception as e_render:
                                                            if log_func_s1: log_func_s1(f"              오류: '{item_name_ocr_log}' 텍스트 렌더링 실패: {e_render}")
                                                    
                                                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                                        progress_callback_item_completed(
                                                            ocr_feedback_location, "ocr_status_rendering_text",
                                                            render_per_block_weight, f"'{item_name_ocr_log}' (블록 {i+1}/{len(translated_ocr_texts)})"
                                                        )
                                                if not (stop_event and stop_event.is_set()): # 중단되지 않았을 때만 남은 가중치 처리
                                                    current_ocr_item_processed_weight += (render_per_block_weight * len(translated_ocr_texts)) # 실제 처리된 가중치
                                            elif ocr_rendering_total_alloc_weight > 0:
                                                if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                                    progress_callback_item_completed(ocr_feedback_location, "ocr_status_rendering_complete_no_text", ocr_rendering_total_alloc_weight, f"'{item_name_ocr_log}'")
                                                current_ocr_item_processed_weight += ocr_rendering_total_alloc_weight
                                            
                                            if stop_event and stop_event.is_set(): break
                                            
                                            if any_text_rendered_on_image:
                                                output_img_stream = io.BytesIO()
                                                save_format = original_img_format if original_img_format and original_img_format.upper() in ['JPEG', 'PNG', 'GIF', 'BMP', 'TIFF'] else 'PNG'
                                                img_to_render_on.save(output_img_stream, format=save_format)
                                                output_img_stream.seek(0)
                                                left, top, width, height = shape_ocr.left, shape_ocr.top, shape_ocr.width, shape_ocr.height
                                                name_orig_img = shape_ocr.name
                                                sp_xml_elem = shape_ocr.element
                                                parent_xml_elem = sp_xml_elem.getparent()
                                                if parent_xml_elem is not None:
                                                    parent_xml_elem.remove(sp_xml_elem)
                                                    new_pic_shape = current_slide_for_ocr.shapes.add_picture(
                                                        output_img_stream, left, top, width=width, height=height
                                                    )
                                                    if name_orig_img: new_pic_shape.name = name_orig_img
                                                    if log_func_s1: log_func_s1(f"        이미지 '{item_name_ocr_log}' (S{slide_idx_ocr+1}) 성공적으로 교체됨.")
                                        else:
                                            if log_func_s1: log_func_s1(f"        경고: '{item_name_ocr_log}' OCR 텍스트 수 불일치. 이미지 변경 없음.")
                                    
                                    if not ocr_texts_for_translation and not (stop_event and stop_event.is_set()): # 번역할 텍스트가 처음부터 없었고 중단되지 않은 경우
                                        remaining_weight_after_detection = ocr_item_total_weight - current_ocr_item_processed_weight
                                        if remaining_weight_after_detection > 0.01 and progress_callback_item_completed:
                                            progress_callback_item_completed(
                                                ocr_feedback_location, "ocr_status_processing_complete_no_translation_needed",
                                                remaining_weight_after_detection, f"'{item_name_ocr_log}'"
                                            )
                                except Exception as e_ocr_img_proc:
                                    logger.error(f"{main_log_prefix} 이미지 OCR 처리 중 예외 '{item_name_ocr_log}': {e_ocr_img_proc}", exc_info=True)
                                    if log_func_s1: log_func_s1(f"      오류: '{item_name_ocr_log}' 이미지 OCR 처리 중 예외: {e_ocr_img_proc}. 건너뜀.")
                                    error_weight_to_report = ocr_item_total_weight - current_ocr_item_processed_weight
                                    if error_weight_to_report > 0.01 and progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                        progress_callback_item_completed(ocr_feedback_location, "ocr_status_processing_error", error_weight_to_report, f"'{item_name_ocr_log}' 오류")
                else: # This corresponds to the original elif condition (image_translation_enabled and not ocr_handler)
                    if log_func_s1: log_func_s1("이미지 번역 활성화되었으나 OCR 핸들러 없어 건너뜁니다.")
            # MODIFIED PART ENDS HERE
            
        except Exception as e_stage1_main: 
            logger.error(f"{main_log_prefix} 1단계 처리 중 심각한 오류: {e_stage1_main}", exc_info=True)
            if log_func_s1: log_func_s1(f"!!! 1단계 처리 중 심각한 오류: {e_stage1_main}\n{traceback.format_exc()}")
            if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
            return False 

        if stop_event and stop_event.is_set():
            if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
            return False

        logger.info(f"{main_log_prefix} 1단계 (텍스트 및 이미지 번역) 성공적 완료.")
        if log_func_s1: log_func_s1(f"--- 1단계: 차트 외 요소 번역 성공적으로 완료 ---\n")
        
        if f_task_log_s1 and not f_task_log_s1.closed:
            try: f_task_log_s1.close()
            except Exception as e_close_log: logger.warning(f"{main_log_prefix} 1단계 작업 로그 파일 닫기 실패: {e_close_log}")
        
        return True